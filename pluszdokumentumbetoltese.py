# imports
from langchain_core.runnables import RunnablePassthrough
from langchain_core.prompts import ChatPromptTemplate
from pydantic import BaseModel, Field
from langchain_core.output_parsers import StrOutputParser
from langchain_neo4j import Neo4jGraph
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_openai import ChatOpenAI
from langchain_experimental.graph_transformers import LLMGraphTransformer
from neo4j import GraphDatabase
from langchain_community.vectorstores import Neo4jVector
from langchain_openai import OpenAIEmbeddings
from langchain.schema import Document
import os
from dotenv import load_dotenv
from langchain_community.document_loaders import PyPDFLoader
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl

load_dotenv()

# Használja ugyanazokat a függvényeket a dokumentum betöltéséhez
def load_pdf(file_path):
    pages = PyPDFLoader(file_path).load()
    documents = []
    for page_num, page in enumerate(pages, start=1):
        paragraphs = page.page_content.split("\n\n")
        for paragraph_num, paragraph in enumerate(paragraphs, start=1):
            documents.append(Document(
                page_content=paragraph.strip(),
                metadata={
                    "source": file_path,
                    "page": page_num,
                    "paragraph": paragraph_num
                }
            ))
    return documents


def load_docx(file_path):
    doc = DocxDocument(file_path)
    text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    return [Document(page_content=text, metadata={"source": file_path})]

def load_text(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        text = f.read()
    return [Document(page_content=text, metadata={"source": file_path})]

def load_pptx(file_path):
    prs = Presentation(file_path)
    text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text')])
    return [Document(page_content=text, metadata={"source": file_path})]

def load_xlsx(file_path):
    wb = openpyxl.load_workbook(file_path)
    text = "\n".join(
        " ".join([str(cell.value) for cell in row if cell.value])
        for sheet in wb.sheetnames for row in wb[sheet].iter_rows()
    )
    return [Document(page_content=text, metadata={"source": file_path})]

def load_file(file_path):
    if file_path.endswith('.pdf'):
        return load_pdf(file_path)
    elif file_path.endswith('.docx'):
        return load_docx(file_path)
    elif file_path.endswith('.txt'):
        return load_text(file_path)
    elif file_path.endswith('.pptx'):
        return load_pptx(file_path)
    elif file_path.endswith('.xlsx'):
        return load_xlsx(file_path)
    else:
        raise ValueError(f"Nem támogatott fájltípus: {file_path}")

# Összefoglaló prompt
summary_prompt = ChatPromptTemplate.from_messages(
    [
        ("system", "You are summarizing documents into a concise, coherent paragraph."),
        ("human", "Summarize the following document: {document_text}"),
    ]
)

# Inicializáljuk a Neo4j gráfot - ez nem törli a meglévő adatokat, csak kapcsolódik
graph = Neo4jGraph()

# Inicializáljuk az LLM-et és a transformer-t
llm = ChatOpenAI(temperature=0, model="gpt-4o-mini")
llm_transformer = LLMGraphTransformer(llm=llm)
summary_chain = summary_prompt | llm

# Dokumentumok összefoglaló funkciója
def summarize_documents(docs):
    """
    Összefoglalja a dokumentumokat és visszatér az összefoglalókkal.
    """
    summaries = []
    for doc in docs:
        summary = summary_chain.invoke({"document_text": doc.page_content})
        summaries.append(f"File: {doc.metadata['source']}\nSummary: {summary}\n")
    return "\n".join(summaries)

# Új dokumentum betöltése - itt adja meg az új dokumentum elérési útját
new_file_path = "/Users/binobenjamin/Documents/Github/tudas2/Technology related blogposts_AI_data driven_automatization_cloud.pdf"  # Módosítsa ezt a tényleges elérési útra

# Új dokumentum betöltése
new_docs = load_file(new_file_path)
print(f"Új dokumentum betöltve: {new_docs}")

# Feldarabolás
text_splitter = RecursiveCharacterTextSplitter(chunk_size=250, chunk_overlap=24)
new_documents = text_splitter.split_documents(documents=new_docs)
print(f"Új feldarabolt dokumentumok: {new_documents}")

# Átalakítás gráf formátumúvá
graph_documents = llm_transformer.convert_to_graph_documents(new_documents)
print(f"Gráf dokumentumok létrehozva")

# Hozzáadás a meglévő gráfhoz
graph.add_graph_documents(
    graph_documents,
    baseEntityLabel=True,
    include_source=True
)
print("Dokumentumok hozzáadva a gráfhoz")

# Vektor index frissítése
vector_index = Neo4jVector.from_existing_graph(
    OpenAIEmbeddings(),
    search_type="hybrid",
    node_label="Document",
    text_node_properties=["text"],
    embedding_node_property="embedding"
)
vector_retriever = vector_index.as_retriever()
print("Vektor index frissítve")

# Új dokumentum összefoglalása
new_summaries = summarize_documents(new_docs)
print(f"Új dokumentum összefoglalója:\n{new_summaries}")

print("Új dokumentum sikeresen hozzáadva az adatbázishoz!")