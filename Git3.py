# imports
from langchain_core.runnables import RunnablePassthrough
from langchain_core.prompts import ChatPromptTemplate
from pydantic import BaseModel, Field
from langchain_core.output_parsers import StrOutputParser
from langchain_neo4j import Neo4jGraph
# from langchain_community.graphs import Neo4jGraph
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_openai import ChatOpenAI
from langchain_community.chat_models import ChatOllama
from langchain_experimental.graph_transformers import LLMGraphTransformer
from neo4j import GraphDatabase
from yfiles_jupyter_graphs import GraphWidget
from langchain_community.vectorstores import Neo4jVector
from langchain_community.document_loaders import TextLoader
from langchain_community.vectorstores.neo4j_vector import remove_lucene_chars
from langchain_ollama import OllamaEmbeddings
from langchain_openai import OpenAIEmbeddings
import os
from langchain_experimental.llms.ollama_functions import OllamaFunctions
from neo4j import Driver
from dotenv import load_dotenv
from langchain_community.document_loaders import PyPDFLoader
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl
from langchain.schema import Document
import gradio as gr # gradio for UI
from dotenv import load_dotenv

load_dotenv()

# Initialize Neo4j graph
graph = Neo4jGraph()

# Load documents of various types 
def load_documents_from_folder(folder_path):
    """
    Betölt minden támogatott dokumentumot a megadott mappából.
    """
    supported_extensions = ['.pdf', '.docx', '.txt', '.pptx', '.xlsx']
    documents = []
    for root, _, files in os.walk(folder_path):
        for file in files:
            if file.startswith("~$"):
                # Ideiglenes fájlok kihagyása
                continue
            if any(file.endswith(ext) for ext in supported_extensions):
                file_path = os.path.join(root, file)
                documents.extend(load_file(file_path))
    return documents

def load_pdf(file_path):
    pages = PyPDFLoader(file_path).load()
    documents = []
    for page_num, page in enumerate(pages, start=1):
        paragraphs = page.page_content.split("\n\n")  # Feltételezzük, hogy bekezdések közé dupla soremelés kerül
        for paragraph_num, paragraph in enumerate(paragraphs, start=1):
            documents.append(Document(
                page_content=paragraph.strip(),
                metadata={
                    "source": file_path,
                    "page": page_num,
                    "paragraph": paragraph_num
                }
            ))
    return documents


def load_docx(file_path):
    doc = DocxDocument(file_path)
    text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
    return [Document(page_content=text, metadata={"source": file_path})]

def load_text(file_path):
    text = TextLoader(file_path).load()
    return [Document(page_content=str(text), metadata={"source": file_path})]

def load_pptx(file_path):
    prs = Presentation(file_path)
    text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text')])
    return [Document(page_content=text, metadata={"source": file_path})]

def load_xlsx(file_path):
    wb = openpyxl.load_workbook(file_path)
    text = "\n".join(
        " ".join([str(cell.value) for cell in row if cell.value])
        for sheet in wb.sheetnames for row in wb[sheet].iter_rows()
    )
    return [Document(page_content=text, metadata={"source": file_path})]

def load_file(file_path):
    if file_path.endswith('.pdf'):
        return load_pdf(file_path)
    elif file_path.endswith('.docx'):
        return load_docx(file_path)
    elif file_path.endswith('.txt'):
        return load_text(file_path)
    elif file_path.endswith('.pptx'):
        return load_pptx(file_path)
    elif file_path.endswith('.xlsx'):
        return load_xlsx(file_path)
    else:
        raise ValueError(f"Nem támogatott fájltípus: {file_path}")
    

    # Specify folder path
folder_path = "/Users/binobenjamin/Documents/Github/tudastarbovites"

# Load all documents from the folder
docs = load_documents_from_folder(folder_path)
print(f"Betöltött dokumentumok: {docs}")

# Check if any documents were loaded
if not docs:
    raise ValueError("Nincsenek betöltött dokumentumok. Ellenőrizd a mappa útvonalát és a fájlformátumokat.")

# Split documents into chunks
text_splitter = RecursiveCharacterTextSplitter(chunk_size=250, chunk_overlap=24)
documents = text_splitter.split_documents(documents=docs)
print(f"Feldarabolt dokumentumok: {documents}")

# Check if any documents were split
if not documents:
    raise ValueError("A dokumentumok feldarabolása nem sikerült.")

# llm2 = OllamaFunctions(model="llama3.1", temperature=0)

# Convert documents to graph format
#llm = OllamaFunctions(model="llama3.2", temperature=0, format="json")

llm = ChatOpenAI(temperature=0, model="gpt-4o-mini")
# llm = ChatOllama(model="llama3.1", temperature=0)
# llm = OllamaFunctions(model="llama3.1", temperature=0)


llm_transformer = LLMGraphTransformer(llm=llm)

graph_documents = llm_transformer.convert_to_graph_documents(documents)
# print(f"Átalakított dokumentumok: {graph_documents}")

# Check if graph_documents is empty
if not graph_documents:
    raise ValueError("A graph_documents lista üres. Ellenőrizd a dokumentum betöltési és feldolgozási folyamatot!")

graph.add_graph_documents(
    graph_documents,
    baseEntityLabel=True,
    include_source=True
)

# embeddings = OllamaEmbeddings(
#    model="mxbai-embed-large",
# )

vector_index = Neo4jVector.from_existing_graph(
    OpenAIEmbeddings(),
    # OllamaEmbeddings(model="llama3.1"),
    search_type="hybrid",
    node_label="Document",
    text_node_properties=["text"],
    embedding_node_property="embedding"
)
vector_retriever = vector_index.as_retriever()

driver = GraphDatabase.driver(
    uri=os.environ["NEO4J_URI"],
    auth=(os.environ["NEO4J_USERNAME"], os.environ["NEO4J_PASSWORD"])
)

# Create fulltext index
def create_fulltext_index(tx):
    query = '''
    CREATE FULLTEXT INDEX `fulltext_entity_id` 
    FOR (n:__Entity__) 
    ON EACH [n.id];
    '''
    tx.run(query)

def create_index():
    with driver.session() as session:
        session.execute_write(create_fulltext_index)
        print("Fulltext index created successfully.")

try:
    create_index()
except Exception as e:
    print(f"Error creating index: {e}")

# Close the driver connection
driver.close()

# Define Entities class
class Entities(BaseModel):
    """Identifying information about entities."""
    
    names: list[str] = Field(
        ..., description="All the person, organization, or business entities that appear in the text"
    )

# Define prompt
prompt = ChatPromptTemplate.from_messages(
    [
        ("system", "You are extracting organization and person entities from the text."),
        ("human", "Use the given format to extract information from the following input: {question}"),
    ]
)

entity_chain = prompt | llm.with_structured_output(Entities)
# entity_chain = prompt | llm2.with_structured_output(Entities)

def generate_full_text_query(input: str) -> str:
    words = [el for el in remove_lucene_chars(input).split() if el]
    if not words:
        return ""
    full_text_query = " AND ".join([f"{word}~2" for word in words])
    print(f"Generated Query: {full_text_query}")
    return full_text_query.strip()

# Fulltext index query
def graph_retriever(question: str) -> str:
    """
    Collects the neighborhood of entities mentioned
    in the question
    """
    result = ""
    entities = entity_chain.invoke(question)
    for entity in entities.names:
        response = graph.query(
            """CALL db.index.fulltext.queryNodes('fulltext_entity_id', $query, {limit:2})
            YIELD node,score
            WITH node
            CALL (node) {
              WITH node
              MATCH (node)-[r:!MENTIONS]->(neighbor)
              RETURN node.id + ' - ' + type(r) + ' -> ' + neighbor.id AS output
              UNION ALL
              WITH node
              MATCH (node)<-[r:!MENTIONS]-(neighbor)
              RETURN neighbor.id + ' - ' + type(r) + ' -> ' +  node.id AS output
            }
            RETURN output LIMIT 50
            """,
            {"query": entity},
        )
        result += "\n".join([el['output'] for el in response])
    return result

def full_retriever(question: str):
    graph_data = graph_retriever(question)
    vector_data = [
        (
            el.page_content,
            el.metadata.get("source", "Ismeretlen forrás"),
            el.metadata.get("page", "N/A"),
            el.metadata.get("paragraph", "N/A")
        )
        for el in vector_retriever.invoke(question)
    ]
    
    vector_context = "\n".join(
        f"#Document: {source} (Page: {page}, Paragraph: {paragraph})\n{content}"
        for content, source, page, paragraph in vector_data
    )
    
    final_data = f"""Graph data:
{graph_data}
Vector data:
{vector_context}
    """
    return final_data, [f"{source} (Page: {page}, Paragraph: {paragraph})" for _, source, page, paragraph in vector_data]


def chat(message, history):
    additional_context, sources = full_retriever(message)
    
    context_message = f"{system_message}\n\nAdditional Context:\n{additional_context}"
    
    messages = [{"role": "system", "content": context_message}] + history + [{"role": "user", "content": message}]
    
    response = openai.chat.completions.create(model=MODEL, messages=messages)
    chat_response = response.choices[0].message.content
    
    if sources:
        source_list = "\n".join(f"- {source}" for source in set(sources))
        chat_response += f"\n\nForrások:\n{source_list}"
    else:
        chat_response += "\n\nForrások: Nincsenek elérhető dokumentumok."
    
    return chat_response




from openai import OpenAI
MODEL = "gpt-4o-mini"
openai = OpenAI()

system_message = "You are a helpful assistant who answers the question based on only the provided context. Limit your answer purely on the context provided, do not add any additional information."
system_message += "After your answer please provide in a list the exact, complete sentences that you used from the context to answer the question."
# system_message += "Give short, courteous answers, no more than 1 sentence. "
system_message += "Always be accurate. If you don't know the answer, say so."
system_message += "Always provide detailed, comprehensive answers, including all relevant information from the context. It should be atleast 5 sentences."

# This function looks rather simpler than the one from my video, because we're taking advantage of the latest Gradio updates



gr.ChatInterface(fn=chat, type="messages").launch(share=True)




summary_prompt = ChatPromptTemplate.from_messages(
    [
        ("system", "You are summarizing documents into a concise, coherent paragraph."),
        ("human", "Summarize the following document: {document_text}"),
    ]
)

summary_chain = summary_prompt | llm

# Dokumentumok összefoglaló funkciója
def summarize_documents(docs):
    """
    Összefoglalja a dokumentumokat és visszatér az összefoglalókkal.
    """
    summaries = []
    for doc in docs:
        summary = summary_chain.invoke({"document_text": doc.page_content})
        summaries.append(f"File: {doc.metadata['source']}\nSummary: {summary}\n")
    return "\n".join(summaries)

# Dokumentumösszefoglalók generálása
summaries = summarize_documents(docs)
print(f"Dokumentumok összefoglalója:\n{summaries}")

# Gradio interfész az összefoglalók megjelenítéséhez
def display_summaries():
    return summaries

gr.Interface(
    fn=display_summaries,
    inputs=None,
    outputs="text",
    title="Dokumentum Összefoglaló"
).launch(share=True)